# Load & đọc PDF, DOCX, PPTX, TXT
from typing import Optional
from PyPDF2 import PdfReader
from pptx import Presentation
import docx

def load_pdf(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PDF."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Lỗi đọc PDF: {e}")
        return None

def load_docx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file DOCX."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Lỗi đọc DOCX: {e}")
        return None

def load_pptx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PPTX."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Lỗi đọc PPTX: {e}")
        return None

def load_txt(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file TXT."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Lỗi đọc TXT: {e}")
        return None

def load_file(file_path: str) -> Optional[str]:
    """Tự động nhận diện và đọc nội dung file PDF, DOCX, PPTX, TXT."""
    if file_path.lower().endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return load_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.lower().endswith(".txt"):
        return load_txt(file_path)
    else:
        print("⚠️ Định dạng file không hỗ trợ.")
        return None