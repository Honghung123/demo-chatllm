# Load & đọc PDF, DOCX, PPTX, TXT
from typing import Optional
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import openpyxl

def load_pdf(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PDF."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return None

def load_docx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file DOCX."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return None

def load_pptx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PPTX."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return None

def load_txt(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file TXT."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return None

def load_xlsx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file XLSX."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            text += "\n"
        return text
    except Exception as e:
        print(f"Error reading XLSX: {e}")
        return None

def load_file(file_path: str) -> Optional[str]:
    """Read file content from PDF, DOCX, PPTX, TXT, XLSX.""" 
    if file_path.lower().endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return load_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.lower().endswith(".txt"):
        return load_txt(file_path)
    elif file_path.lower().endswith(".xlsx"):
        return load_xlsx(file_path)
    else:
        print("Unsupported file format.")
        return None