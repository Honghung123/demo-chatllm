# Load & đọc PDF, DOCX, PPTX, TXT
from typing import Optional
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import os

from utils.file_utils import get_root_path

def load_pdf(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PDF."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return None

def load_docx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file DOCX."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return None

def load_pptx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PPTX."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return None

def load_txt(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file TXT."""
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return None
    
file_directory = f"{get_root_path()}/data/files"

def load_file(file_name: str) -> Optional[str]:
    """Tự động nhận diện và đọc nội dung file PDF, DOCX, PPTX, TXT."""
    file_path = os.path.join(file_directory, file_name)
    if file_path.lower().endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return load_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.lower().endswith(".txt"):
        return load_txt(file_path)
    else:
        print("Unsupported file format.")
        return None