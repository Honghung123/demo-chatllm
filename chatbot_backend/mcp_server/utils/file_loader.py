# Load & đọc PDF, DOCX, PPTX, TXT
from typing import Optional
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import openpyxl

def load_pdf(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PDF."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Error reading PDF: {e}")
        return None

def load_docx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file DOCX."""
    try:
        doc = docx.Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error reading DOCX: {e}")
        return None

def load_pptx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file PPTX."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return None

def load_txt(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file TXT.""" 
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return None

def load_xlsx(file_path: str) -> Optional[str]:
    """Đọc nội dung từ file XLSX."""
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = ""
        for sheet in wb.worksheets:
            text += f"Sheet: {sheet.title}\n"
            for row in sheet.iter_rows(values_only=True):
                text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
            text += "\n"
        return text
    except Exception as e:
        print(f"Error reading XLSX: {e}")
        return None

def load_file(file_path: str) -> Optional[str]:
    """Read file content from PDF, DOCX, PPTX, TXT, XLSX.""" 
    if file_path.lower().endswith(".pdf"):
        return load_pdf(file_path)
    elif file_path.lower().endswith(".docx"):
        return load_docx(file_path)
    elif file_path.lower().endswith(".pptx"):
        return load_pptx(file_path)
    elif file_path.lower().endswith(".txt"):
        return load_txt(file_path)
    elif file_path.lower().endswith(".xlsx"):
        return load_xlsx(file_path)
    else:
        print("Unsupported file format.")
        return None

# # Load & đọc PDF, DOCX, PPTX
# from typing import Optional
# import os

# # Import document processing libraries with error handling
# try:
#     from PyPDF2 import PdfReader
#     PDF_SUPPORT = True
# except ImportError:
#     print("PyPDF2 not installed, PDF support disabled")
#     PDF_SUPPORT = False

# try:
#     from pptx import Presentation
#     PPTX_SUPPORT = True
# except ImportError:
#     print("python-pptx not installed, PPTX support disabled")
#     PPTX_SUPPORT = False

# try:
#     import docx
#     DOCX_SUPPORT = True
# except ImportError:
#     print("python-docx not installed, DOCX support disabled")
#     DOCX_SUPPORT = False

# def load_pdf(file_path: str) -> Optional[str]:
#     """Đọc nội dung từ file PDF."""
#     if not PDF_SUPPORT:
#         print("PyPDF2 not installed, cannot read PDF files")
#         return None
        
#     try:
#         reader = PdfReader(file_path)
#         text = ""
#         for page in reader.pages:
#             text += page.extract_text() or ""
#         return text
#     except Exception as e:
#         print(f"Lỗi đọc PDF: {e}")
#         return None

# def load_docx(file_path: str) -> Optional[str]:
#     """Đọc nội dung từ file DOCX."""
#     if not DOCX_SUPPORT:
#         print("python-docx not installed, cannot read DOCX files")
#         return None
        
#     try:
#         doc = docx.Document(file_path)
#         text = "\n".join([para.text for para in doc.paragraphs])
#         return text
#     except Exception as e:
#         print(f"Lỗi đọc DOCX: {e}")
#         return None

# def load_pptx(file_path: str) -> Optional[str]:
#     """Đọc nội dung từ file PPTX."""
#     if not PPTX_SUPPORT:
#         print("python-pptx not installed, cannot read PPTX files")
#         return None
        
#     try:
#         prs = Presentation(file_path)
#         text = ""
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if hasattr(shape, "text"):
#                     text += shape.text + "\n"
#         return text
#     except Exception as e:
#         print(f"Lỗi đọc PPTX: {e}")
#         return None

# def load_text(file_path: str) -> Optional[str]:
#     """Đọc nội dung từ file text."""
#     try:
#         with open(file_path, 'r', encoding='utf-8') as f:
#             return f.read()
#     except UnicodeDecodeError:
#         try:
#             # Try with a different encoding if UTF-8 fails
#             with open(file_path, 'r', encoding='latin-1') as f:
#                 return f.read()
#         except Exception as e:
#             print(f"Lỗi đọc text file: {e}")
#             return None
#     except Exception as e:
#         print(f"Lỗi đọc text file: {e}")
#         return None

# def load_file(file_path: str) -> Optional[str]:
#     """Tự động nhận diện và đọc nội dung file PDF, DOCX, PPTX, và các file text."""
#     if not os.path.exists(file_path):
#         print(f"⚠️ File không tồn tại: {file_path}")
#         return None
        
#     try:
#         if file_path.lower().endswith(".pdf") and PDF_SUPPORT:
#             return load_pdf(file_path)
#         elif file_path.lower().endswith(".docx") and DOCX_SUPPORT:
#             return load_docx(file_path)
#         elif file_path.lower().endswith(".pptx") and PPTX_SUPPORT:
#             return load_pptx(file_path)
#         elif file_path.lower().endswith((".txt", ".md", ".json", ".py", ".js", ".html", ".css")):
#             return load_text(file_path)
#         else:
#             # Try as text file for unknown file types
#             try:
#                 return load_text(file_path)
#             except:
#                 print("⚠️ Không thể đọc file như file văn bản.")
#                 return None
#     except Exception as e:
#         print(f"⚠️ Lỗi khi đọc file: {e}")
#         return None